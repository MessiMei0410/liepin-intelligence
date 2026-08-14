from __future__ import annotations

import tempfile
import unittest
import zipfile
from datetime import datetime
from pathlib import Path

from a_system_agent.native_attachments import (
    attachment_read_requested,
    detect_wechat_image_bubble,
    image_analysis_requested,
    relevant_visible_attachment_names,
    resolve_wechat_attachments,
    visible_attachment_names,
)


def write_docx(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("word/document.xml", document)


class WeChatAttachmentReadingTest(unittest.TestCase):
    def test_only_explicit_read_requests_open_attachments(self) -> None:
        self.assertTrue(attachment_read_requested("Leo 给我发的简历你看下"))
        self.assertTrue(attachment_read_requested("读取这个 DOCX 的内容"))
        self.assertTrue(attachment_read_requested("准备赵文杰薪资明细与谈薪区间"))
        self.assertTrue(attachment_read_requested("陈明习的简历怎么样"))
        self.assertTrue(attachment_read_requested("这份简历是否合适"))
        self.assertFalse(attachment_read_requested("你能看到这个文件吗"))
        self.assertFalse(attachment_read_requested("总结当前聊天"))

    def test_image_analysis_requires_an_explicit_image_question(self) -> None:
        self.assertTrue(image_analysis_requested("Leo 发的图片你看下"))
        self.assertTrue(image_analysis_requested("识别这张截图里的内容"))
        self.assertFalse(image_analysis_requested("总结当前聊天"))

    def test_image_bubble_detector_prefers_lower_dense_visual_block(self) -> None:
        try:
            import cv2
            import numpy as np
        except ImportError as exc:
            self.skipTest(f"本机未安装 opencv/numpy: {exc}")

        image = np.full((900, 1400, 3), 248, dtype=np.uint8)
        for x in range(980, 1230, 12):
            cv2.line(image, (x, 520), (x, 700), (20, 20, 20), 2)
        for y in range(520, 700, 12):
            cv2.line(image, (980, y), (1230, y), (80, 80, 80), 2)
        ok, encoded = cv2.imencode(".png", image)
        self.assertTrue(ok)
        result = detect_wechat_image_bubble(encoded.tobytes())
        self.assertTrue(result["ok"])
        self.assertGreater(result["box"]["x"], 0.6)
        self.assertGreater(result["box"]["y"], 0.45)

    def test_visible_attachment_names_use_exact_supported_filename_blocks(self) -> None:
        names = visible_attachment_names(
            {
                "text_blocks": [
                    "Leo",
                    "陈明习 20260609B.docx",
                    "33.3K",
                    "不支持.exe",
                ]
            }
        )
        self.assertEqual(names, ["陈明习 20260609B.docx"])

    def test_visible_attachment_names_include_excel(self) -> None:
        names = visible_attachment_names(
            {"text_blocks": ["尹双鹏薪资明细.xls", "候选人薪资结构.xlsx"]}
        )
        self.assertEqual(names, ["尹双鹏薪资明细.xls", "候选人薪资结构.xlsx"])

    def test_relevant_attachment_names_prefer_the_person_named_in_message(self) -> None:
        wechat = {
            "text_blocks": [
                "陈明习 20260609B.docx",
                "尹双鹏薪资明细.xls",
            ]
        }
        self.assertEqual(
            relevant_visible_attachment_names(wechat, "陈明习的简历怎么样"),
            ["陈明习 20260609B.docx"],
        )
        self.assertEqual(
            relevant_visible_attachment_names(wechat, "分析尹双鹏的薪资"),
            ["尹双鹏薪资明细.xls"],
        )
        self.assertEqual(
            relevant_visible_attachment_names(wechat, "这份简历是否合适"),
            ["陈明习 20260609B.docx", "尹双鹏薪资明细.xls"],
        )

    def test_resolver_extracts_xlsx_rows(self) -> None:
        import openpyxl

        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            path = root / "wxid_test/msg/file/2026-07/赵文杰薪资明细.xlsx"
            path.parent.mkdir(parents=True)
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            sheet.append(["项目", "金额", "月数"])
            sheet.append(["固定工资", 45000, 13])
            sheet.append(["年终奖", 45000, 5])
            workbook.save(path)
            evidence = resolve_wechat_attachments(
                {
                    "bridge": {
                        "surface": "native",
                        "wechat": {"text_blocks": ["赵文杰薪资明细.xlsx"]},
                    }
                },
                "准备赵文杰薪资明细与谈薪区间",
                root=root,
            )
            item = evidence["items"][0]
            self.assertTrue(item["content_available"])
            self.assertIn("固定工资\t45000\t13", item["extracted_text"])
            self.assertIn("年终奖\t45000\t5", item["extracted_text"])

    def test_resolver_extracts_pptx_text(self) -> None:
        try:
            from pptx import Presentation
        except ImportError as exc:
            self.skipTest(f"本机未安装 python-pptx: {exc}")

        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            path = root / "wxid_test/msg/file/2026-07/谈薪方案.pptx"
            path.parent.mkdir(parents=True)
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "赵文杰谈薪方案"
            presentation.save(path)
            evidence = resolve_wechat_attachments(
                {
                    "bridge": {
                        "surface": "native",
                        "wechat": {"text_blocks": ["谈薪方案.pptx"]},
                    }
                },
                "打开并总结这个 PPTX",
                root=root,
            )
            self.assertIn("赵文杰谈薪方案", evidence["items"][0]["extracted_text"])

    def test_resolver_prefers_received_file_and_extracts_docx_without_exposing_path(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            received = root / "wxid_test/msg/file/2026-07/陈明习 20260609B.docx"
            dragged = root / "wxid_test/temp/drag/陈明习 20260609B.docx"
            write_docx(received, "8年机械设计经验，负责半导体设备研发。")
            write_docx(dragged, "临时旧副本")
            context = {
                "bridge": {
                    "surface": "native",
                    "wechat": {
                        "text_blocks": ["Leo", "陈明习 20260609B.docx", "33.3K"],
                    },
                }
            }
            evidence = resolve_wechat_attachments(
                context, "Leo 给我发的简历你看下", root=root
            )
            item = evidence["items"][0]
            self.assertTrue(item["content_available"])
            self.assertEqual(item["source_kind"], "wechat_received_file")
            self.assertEqual(item["exact_match_count"], 2)
            self.assertIn("8年机械设计经验", item["extracted_text"])
            self.assertNotIn("临时旧副本", item["extracted_text"])
            self.assertNotIn(str(root), str(evidence))
            self.assertFalse(evidence["chat_database_accessed"])

    def test_resolver_prefers_open_wechat_working_copy(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            root = Path(temp)
            month = datetime.now().strftime("%Y-%m")
            received = root / f"wxid_test/msg/file/{month}/谈薪方案.docx"
            working = root / f"wxid_test/temp/RWTemp/{month}/hash123/谈薪方案.docx"
            write_docx(received, "接收目录旧内容")
            write_docx(working, "当前打开的本地工作副本")
            evidence = resolve_wechat_attachments(
                {
                    "bridge": {
                        "surface": "native",
                        "wechat": {"text_blocks": ["谈薪方案.docx"]},
                    }
                },
                "打开谈薪方案文件",
                root=root,
            )
            item = evidence["items"][0]
            self.assertEqual(item["source_kind"], "wechat_working_copy")
            self.assertIn("当前打开的本地工作副本", item["extracted_text"])
            self.assertNotIn("接收目录旧内容", item["extracted_text"])


if __name__ == "__main__":
    unittest.main()
