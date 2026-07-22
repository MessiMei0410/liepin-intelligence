from __future__ import annotations

import os
import json
import multiprocessing
import re
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree


WECHAT_FILES_ROOT = (
    Path.home()
    / "Library/Containers/com.tencent.xinWeChat/Data/Documents/xwechat_files"
)
WECHAT_ACCOUNT_IDS = tuple(
    value.strip()
    for value in os.environ.get("A_SYSTEM_WECHAT_ACCOUNT_IDS", "").split(",")
    if re.fullmatch(r"wxid_[A-Za-z0-9_]+", value.strip())
)
ATTACHMENT_INDEX_PATH = Path.home() / "Library/Application Support/ASA/attachment-index.json"
SUPPORTED_EXTENSIONS = {".docx", ".pdf", ".txt", ".md", ".csv", ".xls", ".xlsx", ".pptx"}
SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
MAX_ATTACHMENT_BYTES = 25 * 1024 * 1024
MAX_EXTRACTED_CHARS = 18000
_WORD_NAMESPACE = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_ATTACHMENT_PATTERN = re.compile(
    r"([\u4e00-\u9fffA-Za-z0-9][\u4e00-\u9fffA-Za-z0-9 _()（）.\-]{0,159}\.(?:docx|pdf|txt|md|csv|xlsx|xls|pptx))",
    re.IGNORECASE,
)


def attachment_read_requested(message: str) -> bool:
    text = " ".join(str(message or "").split()).lower()
    if not text:
        return False
    object_requested = any(token in text for token in ("简历", "附件", "文件", "文档", "薪资明细", "docx", "pdf", "txt", "md", "csv", "xls", "xlsx", "ppt", "pptx"))
    action_requested = any(token in text for token in (
        "看下", "看看", "查看", "读取", "打开", "分析", "总结", "评估", "匹配", "内容", "准备", "整理",
        "怎么样", "如何", "值不值得", "合不合适", "是否合适",
    ))
    visibility_only = any(token in text for token in ("能看到", "看得到", "是否看到")) and not any(
        token in text for token in ("看下", "看看", "读取", "打开", "分析", "总结", "评估", "匹配", "内容")
    )
    return object_requested and action_requested and not visibility_only


def image_analysis_requested(message: str) -> bool:
    text = " ".join(str(message or "").split()).lower()
    image_requested = any(token in text for token in ("图片", "照片", "截图", "图里", "图中", "这张图"))
    action_requested = any(token in text for token in ("看", "识别", "读取", "分析", "总结", "判断", "内容"))
    return image_requested and action_requested


def detect_wechat_image_bubble(image_data: bytes) -> dict[str, Any]:
    if not image_data or len(image_data) > 16 * 1024 * 1024:
        return {"ok": False, "error": "微信窗口截图为空或超过 16 MB。"}
    try:
        import cv2
        import numpy as np
    except ImportError:
        return {"ok": False, "error": "本机缺少 OpenCV 图片定位能力。"}
    image = cv2.imdecode(np.frombuffer(image_data, dtype=np.uint8), cv2.IMREAD_COLOR)
    if image is None:
        return {"ok": False, "error": "无法解码微信窗口截图。"}
    height, width = image.shape[:2]
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 40, 120)
    tile = max(8, min(height, width) // 120)
    rows, columns = height // tile, width // tile
    mask = np.zeros((rows, columns), dtype=np.uint8)
    for row in range(rows):
        top = row * tile
        if top < height * 0.03 or top > height * 0.82:
            continue
        for column in range(columns):
            left = column * tile
            if left < width * 0.10 or left > width * 0.95:
                continue
            block = edges[top : top + tile, left : left + tile]
            if block.size and float((block > 0).mean()) > 0.07:
                mask[row, column] = 255
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, np.ones((3, 3), dtype=np.uint8), iterations=1)
    count, _, stats, _ = cv2.connectedComponentsWithStats(mask)
    candidates: list[tuple[float, dict[str, float]]] = []
    for index in range(1, count):
        tile_x, tile_y, tile_width, tile_height, component_area = stats[index]
        x, y = int(tile_x * tile), int(tile_y * tile)
        box_width, box_height = int(tile_width * tile), int(tile_height * tile)
        if box_width < width * 0.045 or box_height < height * 0.045:
            continue
        if box_width > width * 0.45 or box_height > height * 0.62:
            continue
        if x < width * 0.10 or x + box_width > width * 0.965:
            continue
        region = edges[y : min(height, y + box_height), x : min(width, x + box_width)]
        edge_density = float((region > 0).mean()) if region.size else 0.0
        if edge_density < 0.085 or int(component_area) < 40:
            continue
        normalized = {
            "x": round(x / width, 6),
            "y": round(y / height, 6),
            "width": round(box_width / width, 6),
            "height": round(box_height / height, 6),
            "edge_density": round(edge_density, 6),
        }
        score = (y / height) * 3.0 + edge_density * 4.0 + min(1.0, component_area / 250.0)
        candidates.append((score, normalized))
    if not candidates:
        return {"ok": False, "error": "未定位到高置信度图片气泡。"}
    candidates.sort(key=lambda item: item[0], reverse=True)
    return {"ok": True, "box": candidates[0][1], "candidate_count": len(candidates)}


def visible_attachment_names(wechat: dict[str, Any]) -> list[str]:
    blocks = wechat.get("text_blocks") if isinstance(wechat.get("text_blocks"), list) else []
    sources = [str(item or "").strip() for item in blocks if str(item or "").strip()]
    if not sources:
        sources = [str(wechat.get("combined_text") or "")]
    names: list[str] = []
    for source in sources:
        for match in _ATTACHMENT_PATTERN.finditer(source):
            name = match.group(1).strip(" .-_[]【】")
            if Path(name).name != name or Path(name).suffix.lower() not in SUPPORTED_EXTENSIONS:
                continue
            if name not in names:
                names.append(name)
    return names[:6]


def relevant_visible_attachment_names(wechat: dict[str, Any], message: str) -> list[str]:
    names = visible_attachment_names(wechat)
    normalized_message = re.sub(r"\s+", "", str(message or "")).lower()
    matched: list[str] = []
    for name in names:
        stem = Path(name).stem
        normalized_name = re.sub(r"\s+", "", name).lower()
        normalized_stem = re.sub(r"\s+", "", stem).lower()
        identity_match = re.match(
            r"^([\u4e00-\u9fff]{2,4})(?=\s|[_\dA-Za-z.\-]|简历|薪资|工资|谈薪|报告|明细)",
            stem,
        )
        identity = identity_match.group(1) if identity_match else ""
        if (
            normalized_name in normalized_message
            or normalized_stem in normalized_message
            or (identity and identity in normalized_message)
        ):
            matched.append(name)
    return matched or names


def visible_wechat_attachment_path(context: dict[str, Any], filename: str) -> Path | None:
    bridge = context if str(context.get("surface") or "").strip().lower() == "native" else {}
    wechat = bridge.get("wechat") if isinstance(bridge.get("wechat"), dict) else {}
    filename = str(filename or "").strip()
    if filename not in visible_attachment_names(wechat):
        return None
    candidates = _wechat_attachment_candidates(WECHAT_FILES_ROOT, filename)
    return candidates[0] if candidates else None


def _working_copy_candidates(root_text: str, account_ids: tuple[str, ...], filename: str, output: Any) -> None:
    root = Path(root_text)
    now = datetime.now()
    found: list[str] = []
    for account_id in account_ids:
        account_dir = root / account_id
        for offset in range(6):
            month_index = now.year * 12 + now.month - 1 - offset
            month_name = f"{month_index // 12:04d}-{month_index % 12 + 1:02d}"
            working_root = account_dir / "temp/RWTemp" / month_name
            if not working_root.is_dir():
                continue
            found.extend(str(path) for path in working_root.glob(f"*/{filename}") if path.is_file())
    output.put(found)


def _cached_working_copy(filename: str) -> Path | None:
    try:
        payload = json.loads(ATTACHMENT_INDEX_PATH.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    raw = str((payload.get("files") or {}).get(filename) or "")
    candidate = Path(raw)
    try:
        candidate.relative_to(WECHAT_FILES_ROOT)
    except (ValueError, OSError):
        return None
    return candidate if candidate.name == filename and candidate.is_file() else None


def _cache_working_copy(filename: str, path: Path) -> None:
    try:
        path.relative_to(WECHAT_FILES_ROOT)
        payload = json.loads(ATTACHMENT_INDEX_PATH.read_text(encoding="utf-8")) if ATTACHMENT_INDEX_PATH.exists() else {}
        files = payload.get("files") if isinstance(payload.get("files"), dict) else {}
        files[filename] = str(path)
        ATTACHMENT_INDEX_PATH.parent.mkdir(parents=True, exist_ok=True)
        ATTACHMENT_INDEX_PATH.write_text(json.dumps({"files": files}, ensure_ascii=False, indent=2), encoding="utf-8")
    except (OSError, ValueError, json.JSONDecodeError):
        return


def _bounded_working_copy_candidates(root: Path, filename: str, timeout_seconds: float = 2.0) -> list[Path]:
    cached = _cached_working_copy(filename)
    if cached:
        return [cached]
    if not WECHAT_ACCOUNT_IDS:
        return []
    context = multiprocessing.get_context("spawn")
    output = context.Queue()
    process = context.Process(
        target=_working_copy_candidates,
        args=(str(root), WECHAT_ACCOUNT_IDS, filename, output),
    )
    process.start()
    process.join(timeout_seconds)
    if process.is_alive():
        process.terminate()
        process.join(1)
        return []
    try:
        paths = [Path(value) for value in output.get_nowait()]
    except Exception:
        paths = []
    if paths:
        _cache_working_copy(filename, paths[0])
    return paths


def _wechat_attachment_candidates(root: Path, filename: str) -> list[Path]:
    if Path(filename).name != filename:
        return []
    candidates: list[Path] = []
    if not root.is_dir():
        return candidates
    if root == WECHAT_FILES_ROOT and WECHAT_ACCOUNT_IDS:
        account_dirs = [root / account_id for account_id in WECHAT_ACCOUNT_IDS]
        working = _bounded_working_copy_candidates(root, filename)
        if working:
            return working
    else:
        account_dirs = list(root.iterdir())
    now = datetime.now()
    month_names = []
    for offset in range(60):
        month_index = now.year * 12 + now.month - 1 - offset
        month_names.append(f"{month_index // 12:04d}-{month_index % 12 + 1:02d}")
    for account_dir in account_dirs:
        if not account_dir.is_dir():
            continue
        if root != WECHAT_FILES_ROOT:
            working_candidates: list[Path] = []
            for month_name in month_names[:6]:
                working_root = account_dir / "temp/RWTemp" / month_name
                if not working_root.is_dir():
                    continue
                working_candidates.extend(
                    candidate
                    for candidate in working_root.glob(f"*/{filename}")
                    if candidate.is_file()
                )
            if working_candidates:
                candidates.extend(working_candidates)
                continue
        received_root = account_dir / "msg/file"
        for month_name in month_names:
            candidate = received_root / month_name / filename
            if candidate.is_file():
                candidates.append(candidate)
        for relative in ("temp/drag", "temp/InputTemp"):
            candidate = account_dir / relative / filename
            if candidate.is_file():
                candidates.append(candidate)
    return candidates


def _extract_docx(path: Path) -> str:
    with zipfile.ZipFile(path) as archive:
        info = archive.getinfo("word/document.xml")
        if info.file_size > 10 * 1024 * 1024:
            raise ValueError("DOCX 正文超过解析上限")
        raw = archive.read(info)
    if b"<!DOCTYPE" in raw or b"<!ENTITY" in raw:
        raise ValueError("DOCX 包含不安全 XML 声明")
    root = ElementTree.fromstring(raw)
    paragraph_tag = f"{{{_WORD_NAMESPACE}}}p"
    text_tag = f"{{{_WORD_NAMESPACE}}}t"
    lines: list[str] = []
    for paragraph in root.iter(paragraph_tag):
        line = "".join(node.text or "" for node in paragraph.iter(text_tag)).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            import fitz
        except ImportError as exc:
            raise ValueError("本机未安装 PDF 文本解析器") from exc
        document = fitz.open(path)
        try:
            return "\n".join(page.get_text("text") for page in list(document)[:60])
        finally:
            document.close()
    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages[:60])


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("本机未安装 PPTX 解析器") from exc
    presentation = Presentation(path)
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        if index > 200:
            break
        lines.append(f"[幻灯片 {index}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and str(text).strip():
                lines.append(str(text).strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [_spreadsheet_cell(cell.text) for cell in row.cells]
                    if any(cells):
                        lines.append("\t".join(cells))
        if sum(len(line) for line in lines) >= MAX_EXTRACTED_CHARS * 2:
            break
    return "\n".join(lines)


def _extract_delimited_text(path: Path) -> str:
    import csv

    raw = path.read_bytes()
    text = ""
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        text = raw.decode("utf-8", errors="replace")
    lines: list[str] = []
    for row in csv.reader(text.splitlines()):
        cells = [_spreadsheet_cell(value) for value in row[:80]]
        if any(cells):
            lines.append("\t".join(cells))
        if sum(len(line) for line in lines) >= MAX_EXTRACTED_CHARS * 2:
            break
    return "\n".join(lines)


def _spreadsheet_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return " ".join(str(value).split())


def _extract_spreadsheet(path: Path) -> str:
    lines: list[str] = []
    if path.suffix.lower() == ".xlsx":
        try:
            import openpyxl
        except ImportError as exc:
            raise ValueError("本机未安装 XLSX 解析器") from exc
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            for sheet in workbook.worksheets[:20]:
                lines.append(f"[工作表] {sheet.title}")
                for row in sheet.iter_rows(max_row=5000, max_col=80, values_only=True):
                    cells = [_spreadsheet_cell(value) for value in row]
                    while cells and not cells[-1]:
                        cells.pop()
                    if any(cells):
                        lines.append("\t".join(cells))
                    if sum(len(line) for line in lines) >= MAX_EXTRACTED_CHARS * 2:
                        return "\n".join(lines)
        finally:
            workbook.close()
        return "\n".join(lines)
    try:
        import xlrd
    except ImportError as exc:
        raise ValueError("本机未安装 XLS 解析器") from exc
    workbook = xlrd.open_workbook(path, on_demand=True)
    try:
        for sheet in workbook.sheets()[:20]:
            lines.append(f"[工作表] {sheet.name}")
            for row_index in range(min(sheet.nrows, 5000)):
                cells = [_spreadsheet_cell(sheet.cell_value(row_index, column)) for column in range(min(sheet.ncols, 80))]
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    lines.append("\t".join(cells))
                if sum(len(line) for line in lines) >= MAX_EXTRACTED_CHARS * 2:
                    return "\n".join(lines)
    finally:
        workbook.release_resources()
    return "\n".join(lines)


def _extract_text(path: Path) -> tuple[str, bool]:
    size = path.stat().st_size
    if size > MAX_ATTACHMENT_BYTES:
        raise ValueError("附件超过 25 MB 读取上限")
    suffix = path.suffix.lower()
    if suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".pdf":
        text = _extract_pdf(path)
    elif suffix in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix == ".csv":
        text = _extract_delimited_text(path)
    elif suffix in {".xls", ".xlsx"}:
        text = _extract_spreadsheet(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    else:
        raise ValueError(f"暂不支持读取 {suffix} 文件")
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return cleaned[:MAX_EXTRACTED_CHARS], len(cleaned) > MAX_EXTRACTED_CHARS


def _extract_text_worker(path_text: str, output: Any) -> None:
    try:
        output.put({"ok": True, "result": _extract_text(Path(path_text))})
    except Exception as exc:
        output.put({"ok": False, "error": str(exc)[:300]})


def _extract_text_bounded(path: Path, timeout_seconds: float = 8.0) -> tuple[str, bool]:
    context = multiprocessing.get_context("spawn")
    output = context.Queue()
    process = context.Process(target=_extract_text_worker, args=(str(path), output))
    process.start()
    process.join(timeout_seconds)
    if process.is_alive():
        process.terminate()
        process.join(1)
        raise ValueError("附件仍在微信云端或正被占用，请先打开下载后重试")
    try:
        payload = output.get_nowait()
    except Exception as exc:
        raise ValueError("附件解析子进程未返回结果") from exc
    if not payload.get("ok"):
        raise ValueError(str(payload.get("error") or "附件解析失败"))
    text, truncated = payload["result"]
    return str(text), bool(truncated)


def extract_local_document(path: Path, timeout_seconds: float = 8.0) -> tuple[str, bool]:
    """Extract a user-selected local document without exposing its path downstream."""
    path = Path(path)
    if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"暂不支持读取 {path.suffix.lower()} 文件")
    return _extract_text_bounded(path, timeout_seconds=timeout_seconds)


def resolve_wechat_attachments(
    context: dict[str, Any], message: str, *, root: Path | None = None
) -> dict[str, Any]:
    if not attachment_read_requested(message):
        return {}
    bridge = context.get("bridge") if isinstance(context.get("bridge"), dict) else {}
    if str(bridge.get("surface") or "").strip().lower() != "native":
        return {}
    wechat = bridge.get("wechat") if isinstance(bridge.get("wechat"), dict) else {}
    names = relevant_visible_attachment_names(wechat, message)
    if not names:
        return {
            "scope": "explicit_visible_filename",
            "chat_database_accessed": False,
            "items": [],
            "status": "当前微信 OCR 未识别到支持读取的附件文件名。",
        }

    root = root or WECHAT_FILES_ROOT
    items: list[dict[str, Any]] = []
    for filename in names[:3]:
        candidates = _wechat_attachment_candidates(root, filename)
        if not candidates:
            items.append(
                {
                    "file_name": filename,
                    "content_available": False,
                    "status": "本机微信文件目录中未找到精确同名附件。",
                }
            )
            continue
        selected = candidates[0]
        try:
            extracted_text, truncated = (
                _extract_text_bounded(selected)
                if root == WECHAT_FILES_ROOT
                else _extract_text(selected)
            )
            items.append(
                {
                    "file_name": filename,
                    "file_type": selected.suffix.lower().lstrip("."),
                    "size_bytes": selected.stat().st_size,
                    "modified_at": datetime.fromtimestamp(selected.stat().st_mtime).isoformat(timespec="seconds"),
                    "source_kind": (
                        "wechat_working_copy"
                        if "/temp/RWTemp/" in selected.as_posix()
                        else "wechat_received_file"
                        if "/msg/file/" in selected.as_posix()
                        else "wechat_temp_file"
                    ),
                    "exact_match_count": len(candidates),
                    "content_available": bool(extracted_text),
                    "extracted_text": extracted_text,
                    "truncated": truncated,
                    "status": "已读取附件正文。" if extracted_text else "附件没有可提取文本。",
                }
            )
        except (OSError, ValueError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
            items.append(
                {
                    "file_name": filename,
                    "content_available": False,
                    "status": f"附件读取失败：{str(exc)[:180]}",
                }
            )
    return {
        "scope": "explicit_visible_filename",
        "chat_database_accessed": False,
        "items": items,
        "status": "已按当前窗口可见文件名检查本机微信附件。",
    }
